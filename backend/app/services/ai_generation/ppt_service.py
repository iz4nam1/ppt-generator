"""
services/ai_generation/ppt_service.py
=======================================
Hackathon-focused PPT generation pipeline.

WHAT MAKES THIS DIFFERENT:
  - Reads actual code files and extracts real technical details
  - Produces both a .pptx AND a .txt intelligence report
  - 3-phase extraction: values → metrics → master brief
  - Tier-based AI routing: critical slides get parallel AI racing
  - Prompts tuned specifically for hackathon judge psychology

OUTPUT:
  generated_presentation.pptx  — judge-ready slides
  project_intelligence.txt      — full intelligence report to reuse anywhere
"""

import asyncio
import copy
import io
import json
import logging
import os
import re
import time
import uuid
from datetime import datetime
from pathlib import Path
from typing import Callable

from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn
from lxml import etree

from app.services.ai_generation.base import AIFeature
from app.services.ai_generation.ai_router import (
    call_by_tier, call_for_planning, call_for_compression,
    call_for_detection, _groq_call, _gemini_call, _cerebras_call,
    HAS_GEMINI, HAS_GROQ, HAS_CEREBRAS,
)
from app.services.ai_generation.prompt_loader import (
    get_persona, get_archetype, get_slide_tier,
    build_slide_prompt, load_personas,
)
from app.config import (
    SITE_NAME, SITE_URL,
    TEXT_EXTENSIONS, IMAGE_EXTENSIONS,
    CHARS_PER_FILE, TEMP_DIR,
)

log = logging.getLogger(__name__)

# ── Load hackathon extraction prompt from file ────────────────────────────────
def _load_extraction_prompt() -> str:
    try:
        path = Path(__file__).parent.parent.parent.parent / "prompts" / "hackathon_extraction.txt"
        return path.read_text(encoding="utf-8")
    except Exception:
        return ""  # falls back to inline prompt

EXTRACTION_PROMPT_TEMPLATE = _load_extraction_prompt()


# ═══════════════════════════════════════════════════════════════════════════════
#  FILE PROCESSING
# ═══════════════════════════════════════════════════════════════════════════════

CODE_EXTENSIONS = {
    ".py", ".js", ".ts", ".jsx", ".tsx", ".java", ".go",
    ".rs", ".cpp", ".c", ".cs", ".rb", ".sh", ".sql",
}
CONFIG_EXTENSIONS = {".json", ".yaml", ".yml", ".toml", ".ini"}
DOC_EXTENSIONS    = {".md", ".txt", ".rst"}


class ContextChunk:
    def __init__(self, label: str, text: str = "", skipped: bool = False):
        self.label   = label
        self.text    = text
        self.skipped = skipped

    @property
    def file_type(self) -> str:
        ext = Path(self.label).suffix.lower()
        if ext in CODE_EXTENSIONS:   return "code"
        if ext in CONFIG_EXTENSIONS: return "config"
        if ext in DOC_EXTENSIONS:    return "docs"
        return "other"


def process_context_file(filename: str, data: bytes) -> ContextChunk:
    ext = Path(filename).suffix.lower()
    if ext in IMAGE_EXTENSIONS:
        return ContextChunk(filename, skipped=True)
    if ext in TEXT_EXTENSIONS or ext in CODE_EXTENSIONS or ext in CONFIG_EXTENSIONS:
        return ContextChunk(filename, data.decode("utf-8", errors="replace"))
    if ext == ".pdf":
        try:
            import pdfplumber
            with pdfplumber.open(io.BytesIO(data)) as pdf:
                text = "\n\n".join(p.extract_text() or "" for p in pdf.pages)
            return ContextChunk(filename, text.strip())
        except Exception:
            return ContextChunk(filename, "[PDF could not be read]")
    try:
        return ContextChunk(filename, data.decode("utf-8", errors="replace"))
    except Exception:
        return ContextChunk(filename, "[Could not read]")


def build_raw_context(chunks: list) -> str:
    return "\n\n".join(
        f"### {c.label}\n{c.text[:CHARS_PER_FILE]}"
        for c in chunks if not c.skipped and c.text.strip()
    )


def categorise_chunks(chunks: list) -> dict:
    cats = {"code": [], "config": [], "docs": [], "other": []}
    for c in chunks:
        if not c.skipped and c.text.strip():
            cats[c.file_type].append(c)
    return cats


def build_files_block(chunks: list, max_chars_each: int = 4000) -> str:
    return "\n\n".join(
        f"### {c.label}\n{c.text[:max_chars_each]}"
        for c in chunks
    )


# ═══════════════════════════════════════════════════════════════════════════════
#  TEMPLATE PARSING — with label detection
# ═══════════════════════════════════════════════════════════════════════════════

LABEL_PATTERNS = re.compile(
    r'^(team\s*name|team\s*leader|problem\s*statement|brief\s*about|'
    r'list\s*of\s*features|process\s*flow|wireframe|architecture\s*diagram|'
    r'technologies\s*utilized|estimated\s*implementation|snapshots|'
    r'prototype\s*performance|additional\s*details|prototype\s*assets|'
    r'your\s*solution\s*should|github\s*public|describe\s*your|'
    r'why\s*ai\s*is\s*required|how\s*aws\s*services|what\s*value)',
    re.IGNORECASE
)


def is_template_label(text: str) -> bool:
    clean = text.strip().rstrip(':').rstrip('?').strip()
    return bool(LABEL_PATTERNS.match(clean))


def find_body_format_para(paras):
    for para in paras[2:]:
        if (para.runs and para.text.strip()
                and not para.runs[0].font.bold
                and not is_template_label(para.text)):
            return para
    for para in paras[1:]:
        if para.runs and para.text.strip() and not is_template_label(para.text):
            return para
    return paras[1] if len(paras) > 1 else paras[0]


def get_slide_info(prs: Presentation) -> list:
    result = []
    for slide_idx, slide in enumerate(prs.slides):
        has_ph, title_text, content_ph = False, None, None
        for ph in slide.placeholders:
            has_ph = True
            idx = ph.placeholder_format.idx
            if idx == 0 and ph.has_text_frame:
                t = ph.text_frame.text.strip()
                if t: title_text = t
            elif idx == 1: content_ph = ph
        if has_ph and title_text:
            result.append({"slide_index": slide_idx, "title": title_text,
                           "mode": "placeholder", "shape": content_ph, "fmt_para": None})
            continue

        text_shapes = sorted(
            [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()],
            key=lambda s: s.top if s.top is not None else 99999
        )
        if not text_shapes: continue

        multi = [s for s in text_shapes if len(s.text_frame.paragraphs) >= 2]
        if multi:
            main  = max(multi, key=lambda s: len(s.text_frame.paragraphs))
            paras = main.text_frame.paragraphs
            title_text = None
            for para in paras:
                if para.text.strip() and not is_template_label(para.text):
                    title_text = para.text.strip()
                    break
            if not title_text:
                title_text = paras[0].text.strip()
            if title_text:
                result.append({"slide_index": slide_idx, "title": title_text,
                               "mode": "textbox_multi", "shape": main,
                               "fmt_para": find_body_format_para(paras)})
            continue

        target     = max(text_shapes, key=lambda s: len(s.text_frame.text))
        title_text = target.text_frame.text.strip()
        if title_text:
            result.append({"slide_index": slide_idx, "title": title_text,
                           "mode": "textbox_single", "shape": target,
                           "fmt_para": None, "all_shapes": text_shapes})
    return result


# ═══════════════════════════════════════════════════════════════════════════════
#  FORMAT-PRESERVING INSERTION
# ═══════════════════════════════════════════════════════════════════════════════

def make_para_xml(text: str, src_para=None, font_size_pt: int = 13) -> etree._Element:
    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    p_el = etree.Element(f"{{{a_ns}}}p")
    if src_para is not None:
        src_pPr = src_para._p.find(qn('a:pPr'))
        if src_pPr is not None: p_el.append(copy.deepcopy(src_pPr))
    rPr_tpl = None
    if src_para is not None and src_para.runs:
        rPr_el = src_para.runs[0]._r.find(qn('a:rPr'))
        if rPr_el is not None: rPr_tpl = copy.deepcopy(rPr_el)
    r_el = etree.SubElement(p_el, f"{{{a_ns}}}r")
    if rPr_tpl is not None:
        r_el.append(rPr_tpl)
    else:
        rPr = etree.SubElement(r_el, f"{{{a_ns}}}rPr")
        rPr.set('lang', 'en-US')
        rPr.set('sz', str(font_size_pt * 100))
        rPr.set('dirty', '0')
    t_el = etree.SubElement(r_el, f"{{{a_ns}}}t")
    t_el.text = text
    return p_el


def clean_lines(raw: str) -> list:
    lines = [l.strip() for l in raw.splitlines() if l.strip()]
    cleaned = []
    for line in lines:
        line = re.sub(r'^[\-\*•▪▸►\d\.\)\|:]+\s*', '', line).strip()
        if is_template_label(line): continue
        if len(line) > 15:
            cleaned.append(line)
    return cleaned


def insert_slide_content(info: dict, lines: list):
    if not lines or info["shape"] is None: return
    shape, mode = info["shape"], info["mode"]

    if mode == "placeholder":
        tf = shape.text_frame; tf.clear()
        sz = None
        try: sz = tf.paragraphs[0].runs[0].font.size
        except: pass
        for line in lines:
            if not line.strip(): continue
            p = tf.add_paragraph(); r = p.add_run()
            r.text = line.strip(); r.font.size = sz or Pt(13)

    elif mode == "textbox_multi":
        tf = shape.text_frame; txBody = tf._txBody
        fmt_para = info.get("fmt_para")
        for p in txBody.findall(qn('a:p'))[1:]:
            txBody.remove(p)
        for line in lines:
            if line.strip():
                txBody.append(make_para_xml(line.strip(), src_para=fmt_para))

    elif mode == "textbox_single":
        tf = shape.text_frame; txBody = tf._txBody
        src = tf.paragraphs[0] if tf.paragraphs else None
        for line in lines:
            if line.strip():
                txBody.append(make_para_xml(line.strip(), src_para=src))


# ═══════════════════════════════════════════════════════════════════════════════
#  PROVENANCE
# ═══════════════════════════════════════════════════════════════════════════════

def embed_provenance(prs: Presentation, generation_id: str):
    try:
        cp = prs.core_properties
        if not cp.author: cp.author = f"Generated by {SITE_NAME}"
        cp.subject  = f"{SITE_NAME}:{generation_id}"
        cp.keywords = f"deckforge hackathon {SITE_URL}"
        cp.comments = (f"Created with {SITE_NAME} ({SITE_URL}) | "
                       f"ID: {generation_id} | {datetime.utcnow().isoformat()}")
    except Exception as e:
        log.warning(f"Provenance embed failed: {e}")


# ═══════════════════════════════════════════════════════════════════════════════
#  PHASE 1: HACKATHON INTELLIGENCE EXTRACTION
#  The core upgrade — extracts like a judge reading your code
# ═══════════════════════════════════════════════════════════════════════════════

async def phase_extract_intelligence(
    description:   str,
    chunks:        list,
    hackathon_info: dict,
    push:          Callable,
) -> tuple[str, str]:
    """
    Extract deep intelligence from the codebase.
    Returns: (master_brief, intelligence_report_text)

    The intelligence_report_text is saved as a .txt file for the user.
    The master_brief is a condensed version used to fill slides.

    3 sub-phases:
      1a. Extract exact values from code    (Cerebras — fast, unlimited)
      1b. Derive metrics from architecture  (Groq — good at reasoning)
      1c. Synthesise master brief           (Gemini — big picture)
    """
    cats      = categorise_chunks(chunks)
    has_code  = bool(cats["code"] or cats["config"])
    has_docs  = bool(cats["docs"] or cats["other"])
    loop      = asyncio.get_event_loop()

    hackathon_name = hackathon_info.get("name", "hackathon")
    hackathon_type = hackathon_info.get("type", "general")
    user_metrics   = hackathon_info.get("metrics", "")
    user_built     = hackathon_info.get("built_vs_planned", "")
    user_user      = hackathon_info.get("who_is_user", "")
    user_gaps      = hackathon_info.get("gaps", "")
    user_demo      = hackathon_info.get("demo_url", "")

    # Build questionnaire block
    q_block = ""
    if any([user_metrics, user_built, user_user, user_gaps, user_demo]):
        q_block = f"""
USER-PROVIDED FACTS (highest priority — use these verbatim):
Real measured metrics: {user_metrics or 'not provided'}
What is built vs planned: {user_built or 'not provided'}
Specific target user: {user_user or 'not provided'}
Known gaps/limitations: {user_gaps or 'not provided'}
Demo URL / GitHub: {user_demo or 'not provided'}
"""

    # ── Sub-phase 1a: Extract exact values ───────────────────────────────────
    push("phase", "Extracting exact values from your code...",
         {"phase": "extract", "step": 1, "total": 3})

    code_block   = build_files_block(cats["code"],   max_chars_each=3500)
    config_block = build_files_block(cats["config"], max_chars_each=2000)

    values_prompt = f"""Extract EVERY hardcoded value, constant, and pattern from this code.

CODE FILES:
{code_block or '(none)'}

CONFIG FILES:
{config_block or '(none)'}

Extract:
1. CONSTANTS & LIMITS: every hardcoded number, string constant, limit, threshold
   e.g. max_tokens=2500, CACHE_VERSION="v3", poll_attempts=30, retry_delay=2s
2. API & SERVICES: every external service, exact API endpoint, SDK method name
   e.g. gemini-2.5-flash, Textract.detect_document_text(), DynamoDB table="civic_ai_cache"
3. DATA STRUCTURES: cache key format, table schemas, request/response structures
   e.g. SHA-256(doc+mode+answers+lang+CACHE_VERSION)
4. ERROR HANDLING: every retry pattern, fallback, exception type caught
   e.g. 3-attempt exponential backoff on HTTP 429, magic bytes validation for PDF/JPEG/PNG
5. INPUT VALIDATION: sanitization methods, size limits, type checks

Quote exact variable names and values. Be exhaustive."""

    values_result = await loop.run_in_executor(
        None, _cerebras_call, values_prompt, "", 1500, "extract-values"
    )
    if not values_result:
        values_result = await loop.run_in_executor(
            None, _groq_call, values_prompt, "",
            "llama-3.3-70b-versatile", 1500, False, "extract-values-fb"
        )
    values_result = values_result or "No code files provided."
    await asyncio.sleep(1)

    # ── Sub-phase 1b: Derive metrics ─────────────────────────────────────────
    push("phase", "Deriving performance and cost metrics...",
         {"phase": "extract", "step": 2, "total": 3})

    metrics_prompt = f"""You are a senior engineer estimating system performance.

PROJECT: {description}
HACKATHON: {hackathon_name}

EXTRACTED CODE VALUES:
{values_result[:2000]}

CODE ARCHITECTURE (main files):
{build_files_block(cats['code'][:2], max_chars_each=2500) or '(none)'}

{q_block}

Derive realistic performance and cost estimates.
For each estimate, show your reasoning briefly.

LATENCY ESTIMATES:
For each operation, estimate based on services used. Label as [from code] or [estimated from architecture].
Format: operation → estimate [reasoning]

COST ESTIMATES at 1,000 users/month:
Calculate from actual API calls found in code × published pricing.
Show: service → monthly operations → unit cost → monthly total
Label: [calculated: X × $Y]

BOTTLENECKS:
What is the slowest step? What has no caching? What fails under load?

RELIABILITY:
What retry logic exists? What has no error handling? Single points of failure?"""

    metrics_result = await loop.run_in_executor(
        None, _groq_call, metrics_prompt, "",
        "llama-3.3-70b-versatile", 1500, False, "derive-metrics"
    )
    if not metrics_result:
        metrics_result = await loop.run_in_executor(
            None, _gemini_call, metrics_prompt, "", False, "derive-metrics-fb"
        )
    metrics_result = metrics_result or "Metrics could not be derived."
    await asyncio.sleep(1)

    # ── Sub-phase 1c: Full intelligence report ────────────────────────────────
    push("phase", "Building complete intelligence report...",
         {"phase": "extract", "step": 3, "total": 3})

    docs_block = build_files_block(cats["docs"] + cats["other"], max_chars_each=3000)

    # Use the external prompt template if loaded, else use inline
    if EXTRACTION_PROMPT_TEMPLATE:
        full_report_prompt = f"""{EXTRACTION_PROMPT_TEMPLATE}

PROJECT: {description}
HACKATHON: {hackathon_name} ({hackathon_type})

DOCUMENTATION & README:
{docs_block or '(none provided)'}

EXTRACTED CODE VALUES:
{values_result[:2000]}

DERIVED METRICS:
{metrics_result[:1500]}

{q_block}

Write the complete intelligence report following the template above.
Every fact must be traceable to the inputs. Do not invent numbers not in the inputs.
Label all estimates clearly."""
    else:
        full_report_prompt = f"""Write a complete hackathon intelligence report for this project.

PROJECT: {description}
HACKATHON: {hackathon_name}

DOCUMENTATION:
{docs_block or '(none)'}

EXTRACTED VALUES:
{values_result[:2000]}

METRICS:
{metrics_result[:1500]}

{q_block}

Include: problem evidence, exact technical details, performance, costs, features,
gaps, slide-ready facts, and a 200-word universal AI prompt."""

    intelligence_report = await loop.run_in_executor(
        None, _gemini_call, full_report_prompt, "", False, "full-report"
    )
    if not intelligence_report:
        intelligence_report = await loop.run_in_executor(
            None, _groq_call, full_report_prompt, "",
            "llama-3.3-70b-versatile", 4096, False, "full-report-fb"
        )
    intelligence_report = intelligence_report or f"# {description}\n\nIntelligence extraction failed."

    # ── Build the condensed brief for slide filling ───────────────────────────
    brief_prompt = f"""Condense this intelligence report into a master brief for filling hackathon slides.

INTELLIGENCE REPORT:
{intelligence_report[:4000]}

{q_block}

Write a PROJECT MASTER BRIEF with these sections:
- CORE DESCRIPTION (2 sentences)
- PROBLEM FACTS (specific stats and named populations)
- SOLUTION MECHANISM (exact step-by-step)
- KEY TECHNOLOGIES (name + version + role + why)
- VERIFIED METRICS (every number, labelled as from code/estimated/measured)
- FEATURES (name + mechanism + benefit)
- ARCHITECTURE (every component named)
- HONEST GAPS (what isn't built)
- SLIDE FACTS (top 15 specific verifiable facts)

Keep it dense. Every sentence must contain something specific."""

    master_brief = await loop.run_in_executor(
        None, _groq_call, brief_prompt, "",
        "llama-3.3-70b-versatile", 2000, False, "master-brief"
    )
    if not master_brief:
        # Fall back to using the full report as brief
        master_brief = intelligence_report[:3000]

    return master_brief, intelligence_report


# ═══════════════════════════════════════════════════════════════════════════════
#  PHASE 2: PLAN SLIDES
# ═══════════════════════════════════════════════════════════════════════════════

async def phase_plan(
    brief:         str,
    slides_info:   list,
    persona:       dict,
    hackathon_info: dict,
) -> dict:
    """Plan content distribution across all slides."""
    titles   = "\n".join(f"  Slide {i+1}: {s['title']}" for i, s in enumerate(slides_info))
    hack_ctx = hackathon_info.get("name", "hackathon")

    prompt = f"""Plan content for a {len(slides_info)}-slide {hack_ctx} presentation.

MASTER BRIEF:
{brief}

SLIDES:
{titles}

AUDIENCE: {persona['judge_profile']}

HACKATHON PLANNING RULES:
1. No fact appears on more than ONE slide
2. Every key_point references a SPECIFIC detail from the brief
3. If brief says "NOT PROVIDED" — do not plan to use it
4. Problem slide must have evidence, not assertions
5. Architecture slide must name every component
6. Cost slide must have actual numbers
7. The deck tells a coherent story from problem → solution → proof

Return ONLY valid JSON:
{{"1":{{"angle":"...","key_points":["specific fact 1","fact 2","fact 3","fact 4"],"avoid":"...","opening":"..."}},"2":{{...}}}}"""

    loop   = asyncio.get_event_loop()
    result = await loop.run_in_executor(
        None, _gemini_call, prompt, "", True, "plan"
    )
    if not result:
        result = await loop.run_in_executor(
            None, _groq_call, prompt, "",
            "llama-3.3-70b-versatile", 2000, True, "plan-fb"
        )

    raw = re.sub(r'^```[a-z]*\n?', '', result or '', flags=re.IGNORECASE)
    raw = re.sub(r'\n?```$', '', raw)
    match = re.search(r'\{.*\}', raw, re.DOTALL)
    if match: raw = match.group()

    try:
        return {str(k): v for k, v in json.loads(raw).items()}
    except Exception:
        log.warning("Slide plan JSON parse failed")
        return {}


# ═══════════════════════════════════════════════════════════════════════════════
#  PHASE 3: GENERATE SLIDES
# ═══════════════════════════════════════════════════════════════════════════════

async def phase_generate_slide(
    info:          dict,
    slide_number:  int,
    total:         int,
    brief:         str,
    slide_plan:    dict,
    persona:       dict,
) -> list:
    """Generate content for one slide using tier-based routing."""
    slide_title = info["title"]
    tier        = get_slide_tier(slide_title)
    system, user = build_slide_prompt(
        slide_title  = slide_title,
        slide_number = slide_number,
        total_slides = total,
        brief        = brief,
        slide_plan   = slide_plan,
        persona      = persona,
    )

    raw = await call_by_tier(
        tier   = tier,
        prompt = user,
        system = system,
        label  = f"slide{slide_number}-t{tier}",
    )

    lines = clean_lines(raw)
    return lines if lines else [f"Content for: {slide_title}"]


# ═══════════════════════════════════════════════════════════════════════════════
#  PPT FEATURE CLASS
# ═══════════════════════════════════════════════════════════════════════════════

class PPTGenerationService(AIFeature):
    """
    Hackathon PPT generation.
    Outputs: .pptx file + .txt intelligence report
    """
    FEATURE_KEY = "ppt_generation"

    async def run(
        self,
        template_path:  str,
        description:    str,
        chunks:         list,
        generation_id:  str,
        push:           Callable,
        user_plan:      str  = "free",
        hackathon_info: dict = None,
    ):
        hackathon_info = hackathon_info or {}

        providers = [p for p, v in [
            ("Gemini", HAS_GEMINI),
            ("Groq", HAS_GROQ),
            ("Cerebras", HAS_CEREBRAS)
        ] if v]
        log.info(f"[{generation_id}] Providers: {providers}")

        # ── Parse template ────────────────────────────────────────────────────
        push("progress", "Parsing your template...", {"step": 0, "total_steps": 4})
        try:
            prs         = Presentation(template_path)
            slides_info = get_slide_info(prs)
        except Exception as e:
            push("error", "Could not read the template. Please check it is a valid .pptx.")
            return None, None

        if not slides_info:
            push("error", "No fillable slides found in the template.")
            return None, None

        total_slides = len(slides_info)
        push("info", f"Found {total_slides} slides — using {', '.join(providers)}")

        # ── Detect project type ───────────────────────────────────────────────
        push("progress", "Detecting hackathon category...", {"step": 1, "total_steps": 4})
        raw_context = build_raw_context(chunks)
        raw_type    = await asyncio.get_event_loop().run_in_executor(
            None, call_for_detection, description, raw_context
        )
        clean_type = re.sub(r'[^a-z_]', '', raw_type.lower().strip())
        personas   = load_personas()
        persona    = personas.get(clean_type, personas["technical"])
        push("info", f"Category: {persona['label']}")
        await asyncio.sleep(1)

        # ── Phase 1: Intelligence extraction ─────────────────────────────────
        push("progress", "Extracting project intelligence...", {"step": 2, "total_steps": 4})
        try:
            master_brief, intelligence_report = await phase_extract_intelligence(
                description    = description,
                chunks         = chunks,
                hackathon_info = hackathon_info,
                push           = push,
            )
        except Exception as e:
            log.error(f"Intelligence extraction failed: {e}")
            push("error", "AI service unavailable. Please try again.")
            return None, None

        # ── Phase 2: Plan ─────────────────────────────────────────────────────
        push("progress", "Planning slide content...", {"step": 3, "total_steps": 4})
        try:
            slide_plan = await phase_plan(master_brief, slides_info, persona, hackathon_info)
        except Exception as e:
            log.warning(f"Planning failed (non-fatal): {e}")
            slide_plan = {}

        # ── Phase 3: Generate slides ──────────────────────────────────────────
        push("progress", f"Generating {total_slides} slides...", {"step": 4, "total_steps": 4})

        for i, info in enumerate(slides_info):
            slide_title = info["title"]
            tier        = get_slide_tier(slide_title)
            tier_labels = {1: "⚡ parallel race", 2: "fast", 3: "turbo"}
            short_title = slide_title[:42] + "..." if len(slide_title) > 42 else slide_title

            push("slide_start",
                 f"Slide {i+1}/{total_slides}: {short_title}",
                 {"slide": i+1, "total": total_slides,
                  "title": short_title, "tier": tier,
                  "tier_label": tier_labels.get(tier, "")})
            try:
                lines = await phase_generate_slide(
                    info         = info,
                    slide_number = i + 1,
                    total        = total_slides,
                    brief        = master_brief,
                    slide_plan   = slide_plan,
                    persona      = persona,
                )
                insert_slide_content(info, lines)
                push("slide_done", f"Slide {i+1} done",
                     {"slide": i+1, "total": total_slides, "tier": tier})
            except Exception as e:
                log.error(f"Slide {i+1} failed: {e}")
                push("slide_error", f"Slide {i+1} had an issue.",
                     {"slide": i+1})

            await asyncio.sleep(0.5 if tier == 1 else 1.0)

        # ── Save PPT ──────────────────────────────────────────────────────────
        embed_provenance(prs, generation_id)
        pptx_path = TEMP_DIR / f"{generation_id}_output.pptx"
        prs.save(str(pptx_path))

        # ── Save intelligence report as .txt ──────────────────────────────────
        txt_path = TEMP_DIR / f"{generation_id}_intelligence.txt"
        report_with_header = (
            f"DECKFORGE INTELLIGENCE REPORT\n"
            f"{'='*50}\n"
            f"Project: {description[:100]}\n"
            f"Hackathon: {hackathon_info.get('name', 'Not specified')}\n"
            f"Generated: {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}\n"
            f"Category: {persona['label']}\n"
            f"{'='*50}\n\n"
            f"HOW TO USE THIS REPORT:\n"
            f"- Paste the 'PASTE INTO ANY AI' section into Claude/ChatGPT for any task\n"
            f"- Use 'SLIDE-BY-SLIDE AMMUNITION' when editing individual slides\n"
            f"- Use 'HACKATHON-READY FACTS' for social media, README, pitch emails\n"
            f"{'='*50}\n\n"
            f"{intelligence_report}"
        )
        txt_path.write_text(report_with_header, encoding="utf-8")

        push("done", "Both files ready!",
             {"job_id": generation_id, "has_report": True})

        return str(pptx_path), str(txt_path)
